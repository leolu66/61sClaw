"""
多页 HTML 演示文稿 -> PPTX 转换器

支持三种格式自动检测：
  A) html-ppt-skill:   .deck > section.slide, 用 .is-active 类切换 (opacity)
  B) guizang-ppt:      #deck > .slide, 用 transform:translateX 翻页
  C) 单页:             slideflow 原生单页 HTML

每页处理：截图背景 + DOM 提取文本坐标 -> python-pptx 矢量重建

用法:
    python html_deck_to_pptx.py <html_dir> [output.pptx]
"""

import os, sys, asyncio
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "repos" / "SlideFlow"))

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FONT_MAPPING = {
    "Noto Sans SC": "微软雅黑", "Noto Serif SC": "SimSun",
    "Playfair Display": "Georgia", "Source Serif 4": "Georgia",
    "IBM Plex Mono": "Consolas", "JetBrains Mono": "Consolas",
    "Poppins": "Arial", "Inter": "Segoe UI", "Roboto": "Arial",
}

def system_font(ff):
    primary = ff.split(',')[0].strip('"').strip("'").strip()
    return FONT_MAPPING.get(primary, primary)


async def detect_format(page):
    """Detect the HTML presentation format. Returns (format, slide_count)."""
    result = await page.evaluate("""
        (() => {
            // Format D: huashu-design (#deck > .slide-page, display block/none)
            const hs = document.querySelectorAll('#deck > .slide-page, .deck > .slide-page');
            if (hs.length > 0) return {fmt: 'huashu', count: hs.length};
            
            // Format A: html-ppt-skill (.deck > .slide, section-based)
            const hts = document.querySelectorAll('.deck .slide');
            if (hts.length > 0) return {fmt: 'html-ppt', count: hts.length};
            
            // Format B: guizang-ppt (#deck > .slide, transform-based)
            const gz = document.querySelectorAll('#deck .slide');
            if (gz.length > 0) return {fmt: 'guizang', count: gz.length};
            
            // Format C: single page
            return {fmt: 'single', count: 1};
        })()
    """)
    return result['fmt'], result['count']


async def navigate_slide(page, fmt, idx):
    """Navigate to slide index using format-specific method."""
    if fmt == 'huashu':
        await page.evaluate(f"""
            (() => {{
                const slides = document.querySelectorAll('#deck > .slide-page, .deck > .slide-page');
                for (let i = 0; i < slides.length; i++) slides[i].style.display = i === {idx} ? 'block' : 'none';
            }})()
        """)
    elif fmt == 'html-ppt':
        await page.evaluate(f"""
            (() => {{
                document.querySelectorAll('.slide').forEach((s,i) => {{
                    s.classList.toggle('is-active', i === {idx});
                }});
            }})()
        """)
    elif fmt == 'guizang':
        await page.evaluate(f"""
            (() => {{
                const deck = document.getElementById('deck');
                if (deck) deck.style.transform = 'translateX({-idx*100}vw)';
            }})()
        """)
    # fmt == 'single': no navigation needed
    await page.wait_for_timeout(300)


async def extract_text_elements(page, fmt):
    """Extract text elements from the currently active/visible slide."""
    if fmt == 'huashu':
        js = """
        (() => {
            const els = [];
            const slides = document.querySelectorAll('#deck > .slide-page, .deck > .slide-page');
            let currentSlide = null;
            for (const s of slides) { if (s.style.display !== 'none' && s.style.display !== '') { currentSlide = s; break; } }
            if (!currentSlide) currentSlide = slides[0];
            if (!currentSlide) return els;
            const slideRect = currentSlide.getBoundingClientRect();
            currentSlide.querySelectorAll('*:not(script):not(style):not(svg):not(canvas)').forEach(el => {
                const s = window.getComputedStyle(el);
                if (s.display === 'none' || s.visibility === 'hidden' || s.opacity === '0') return;
                if (el.tagName === 'IMG' || el.tagName === 'svg' || el.tagName === 'CANVAS') return;
                let hasText = false;
                for (const node of el.childNodes) {
                    if (node.nodeType === 3 && node.textContent.trim()) { hasText = true; break; }
                }
                if (!hasText) return;
                const color = s.color.match(/rgb\\((\\d+),\\s*(\\d+),\\s*(\\d+)\\)/);
                if (!color) return;
                const rgb = [parseInt(color[1]), parseInt(color[2]), parseInt(color[3])];
                let text = '';
                for (const node of el.childNodes) {
                    if (node.nodeType === 3) text += node.textContent;
                }
                text = text.trim();
                if (!text || text.length > 500) text = text.slice(0, 500);
                if (!text) return;
                const r = el.getBoundingClientRect();
                els.push({text:text,x:r.left-slideRect.left,y:r.top-slideRect.top,w:r.width,h:r.height,fs:parseFloat(s.fontSize),fw:s.fontWeight,ff:s.fontFamily,color:rgb,ta:s.textAlign});
            });
            return els;
        })()
        """
    elif fmt == 'html-ppt':
        js = """
        (() => {
            const els = [];
            // Only extract from .is-active slide
            const slide = document.querySelector('.slide.is-active');
            if (!slide) return els;
            const slideRect = slide.getBoundingClientRect();
            
            slide.querySelectorAll('*:not(script):not(style):not(svg):not(canvas)').forEach(el => {
                const s = window.getComputedStyle(el);
                if (s.display === 'none' || s.visibility === 'hidden' || s.opacity === '0') return;
                if (el.tagName === 'IMG' || el.tagName === 'svg' || el.tagName === 'CANVAS') return;
                
                let hasText = false;
                for (const node of el.childNodes) {
                    if (node.nodeType === 3 && node.textContent.trim()) { hasText = true; break; }
                }
                if (!hasText) return;
                
                const color = s.color.match(/rgb\\((\\d+),\\s*(\\d+),\\s*(\\d+)\\)/);
                if (!color) return;
                const rgb = [parseInt(color[1]), parseInt(color[2]), parseInt(color[3])];
                
                let text = '';
                for (const node of el.childNodes) {
                    if (node.nodeType === 3) text += node.textContent;
                }
                text = text.trim();
                if (!text || text.length > 500) text = text.slice(0, 500);
                if (!text) return;
                
                const r = el.getBoundingClientRect();
                els.push({
                    text: text,
                    x: r.left - slideRect.left,
                    y: r.top - slideRect.top,
                    w: r.width,
                    h: r.height,
                    fs: parseFloat(s.fontSize),
                    fw: s.fontWeight,
                    ff: s.fontFamily,
                    color: rgb,
                    ta: s.textAlign
                });
            });
            return els;
        })()
        """
    elif fmt == 'guizang':
        js = """
        (() => {
            const els = [];
            const deck = document.getElementById('deck');
            const slides = deck ? deck.querySelectorAll('.slide') : [];
            let idx = 0;
            if (deck && deck.style.transform) {
                const m = deck.style.transform.match(/translateX\\((-?\\d+)/);
                if (m) idx = Math.abs(parseInt(m[1])) / 100;
            }
            const currentSlide = slides[idx] || document.body;
            const slideRect = currentSlide.getBoundingClientRect();
            
            currentSlide.querySelectorAll('*:not(script):not(style):not(canvas)').forEach(el => {
                const s = window.getComputedStyle(el);
                if (s.display === 'none' || s.visibility === 'hidden' || s.opacity === '0') return;
                if (el.tagName === 'IMG' || el.tagName === 'svg' || el.tagName === 'CANVAS') return;
                
                let hasText = false;
                for (const node of el.childNodes) {
                    if (node.nodeType === 3 && node.textContent.trim()) { hasText = true; break; }
                }
                if (!hasText) return;
                
                const color = s.color.match(/rgb\\((\\d+),\\s*(\\d+),\\s*(\\d+)\\)/);
                if (!color) return;
                const rgb = [parseInt(color[1]), parseInt(color[2]), parseInt(color[3])];
                
                let text = '';
                for (const node of el.childNodes) {
                    if (node.nodeType === 3) text += node.textContent;
                }
                text = text.trim();
                if (!text || text.length > 500) text = text.slice(0, 500);
                if (!text) return;
                
                const r = el.getBoundingClientRect();
                els.push({
                    text: text,
                    x: r.left - slideRect.left,
                    y: r.top - slideRect.top,
                    w: r.width,
                    h: r.height,
                    fs: parseFloat(s.fontSize),
                    fw: s.fontWeight,
                    ff: s.fontFamily,
                    color: rgb,
                    ta: s.textAlign
                });
            });
            return els;
        })()
        """
    else:  # single
        js = """
        (() => {
            const els = [];
            const slideRect = document.body.getBoundingClientRect();
            document.querySelectorAll('body *:not(script):not(style):not(canvas)').forEach(el => {
                const s = window.getComputedStyle(el);
                if (s.display === 'none' || s.visibility === 'hidden' || s.opacity === '0') return;
                if (el.tagName === 'IMG' || el.tagName === 'svg' || el.tagName === 'CANVAS') return;
                
                let hasText = false;
                for (const node of el.childNodes) {
                    if (node.nodeType === 3 && node.textContent.trim()) { hasText = true; break; }
                }
                if (!hasText) return;
                
                const color = s.color.match(/rgb\\((\\d+),\\s*(\\d+),\\s*(\\d+)\\)/);
                if (!color) return;
                const rgb = [parseInt(color[1]), parseInt(color[2]), parseInt(color[3])];
                
                let text = '';
                for (const node of el.childNodes) {
                    if (node.nodeType === 3) text += node.textContent;
                }
                text = text.trim();
                if (!text || text.length > 500) text = text.slice(0, 500);
                if (!text) return;
                
                const r = el.getBoundingClientRect();
                els.push({
                    text: text,
                    x: r.left - slideRect.left,
                    y: r.top - slideRect.top,
                    w: r.width,
                    h: r.height,
                    fs: parseFloat(s.fontSize),
                    fw: s.fontWeight,
                    ff: s.fontFamily,
                    color: rgb,
                    ta: s.textAlign
                });
            });
            return els;
        })()
        """
    return await page.evaluate(js)


async def convert_deck(html_dir, ppt_path):
    if not os.path.isdir(html_dir):
        print(f"Error: not a directory: {html_dir}")
        return False

    files = sorted([f for f in os.listdir(html_dir) if f.endswith('.html')])
    if not files:
        print("Error: no HTML files found")
        return False

    html_file = os.path.abspath(os.path.join(html_dir, files[0]))
    print(f"Source: {html_file}")

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(viewport={"width": 1280, "height": 720}, device_scale_factor=2)
        page = await context.new_page()

        try:
            await page.goto(f"file://{html_file}", wait_until="domcontentloaded", timeout=15000)
        except Exception:
            print("  Page load timeout, proceeding...")
        await page.wait_for_timeout(2000)

        # 注入关键 CSS（外部 CSS 路径可能解析失败）
        await page.add_style_tag(content="""
            .deck { position:relative; width:100vw; height:100vh; overflow:hidden; }
            .slide { position:absolute; inset:0; opacity:0; pointer-events:none; overflow:hidden;
                     display:flex; flex-direction:column; justify-content:center;
                     padding:72px 96px; box-sizing:border-box; }
            .slide.is-active { opacity:1; pointer-events:auto; z-index:2; }
        """)
        await page.wait_for_timeout(500)

        fmt, total = await detect_format(page)

        # If CSS is external, the page might show all slides stacked. For html-ppt format, activate slide 0.
        if fmt == 'html-ppt':
            await navigate_slide(page, fmt, 0)

        print(f"Format: {fmt} | Slides: {total}")
        print(f"Output: {ppt_path}")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i in range(total):
            await navigate_slide(page, fmt, i)

            # 1) Hide text, take background screenshot
            await page.evaluate("""
                (()=>{const s=document.createElement('style');s.id='_bg_hide';
                s.textContent='*{color:transparent!important;text-shadow:none!important;-webkit-text-stroke:0!important}';
                document.head.appendChild(s);})()
            """)
            await page.wait_for_timeout(100)
            bg_path = f"_bg_{i}.png"
            await page.screenshot(path=bg_path)

            # 2) Restore text, extract elements
            await page.evaluate("(()=>{const s=document.getElementById('_bg_hide');if(s)s.remove()})()")
            await page.wait_for_timeout(100)

            elements = await extract_text_elements(page, fmt)

            # 3) Build slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(bg_path)

            cnt = 0
            for el in elements:
                x, y = el['x'] * 0.75, el['y'] * 0.75
                w, h = el['w'] * 0.75, el['h'] * 0.75
                if w <= 0 or h <= 0: continue

                tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
                tf = tb.text_frame
                tf.word_wrap = False
                tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0

                p = tf.paragraphs[0]
                ta = el.get('ta', 'left')
                if ta == 'center': p.alignment = PP_ALIGN.CENTER
                elif ta == 'right': p.alignment = PP_ALIGN.RIGHT

                run = p.add_run()
                run.text = el['text']
                run.font.size = Pt(el['fs'] * 0.75)
                run.font.color.rgb = RGBColor(*el['color'])
                run.font.name = system_font(el['ff'])
                fw = el.get('fw', 'normal')
                if fw == 'bold' or (fw.isdigit() and int(fw) >= 600):
                    run.font.bold = True
                cnt += 1

            print(f"  [{i+1}/{total}] {cnt} text elements")

        await browser.close()

    prs.save(ppt_path)
    print(f"Done: {ppt_path} ({os.path.getsize(ppt_path)/1024:.1f} KB)")
    return True


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python html_deck_to_pptx.py <html_dir> [output.pptx]")
        sys.exit(1)
    html_dir = sys.argv[1]
    ppt_path = sys.argv[2] if len(sys.argv) > 2 else os.path.join(html_dir, "output.pptx")
    asyncio.run(convert_deck(html_dir, ppt_path))
