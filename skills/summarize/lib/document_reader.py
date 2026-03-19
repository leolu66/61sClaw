#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档读取器 - 统一接口读取各种文档格式
整合现有技能的功能，无需重复开发
"""

import os
import sys

class DocumentReader:
    """统一文档读取器"""
    
    # 支持的文件格式
    SUPPORTED_FORMATS = {
        '.txt': 'text',
        '.md': 'text',
        '.py': 'text',
        '.js': 'text',
        '.html': 'text',
        '.json': 'text',
        '.xml': 'text',
        '.csv': 'text',
        '.pdf': 'pdf',
        '.docx': 'word',
        '.doc': 'word',
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint',
    }
    
    @classmethod
    def can_read(cls, file_path):
        """检查是否支持该文件格式"""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in cls.SUPPORTED_FORMATS
    
    @classmethod
    def read(cls, file_path, **options):
        """
        读取文档内容
        
        Args:
            file_path: 文件路径
            **options: 额外选项
                - pdf_pages: PDF 页面范围 "1-5" 或 "1,3,5"
                - excel_sheet: Excel 工作表名
                - pptx_slide: PPT 幻灯片范围
        
        Returns:
            dict: 读取结果
        """
        if not os.path.exists(file_path):
            return {
                'success': False,
                'error': '文件不存在'
            }
        
        ext = os.path.splitext(file_path)[1].lower()
        doc_type = cls.SUPPORTED_FORMATS.get(ext)
        
        if not doc_type:
            return {
                'success': False,
                'error': f'不支持的文件格式: {ext}'
            }
        
        try:
            if doc_type == 'text':
                return cls._read_text(file_path)
            elif doc_type == 'pdf':
                return cls._read_pdf(file_path, options.get('pdf_pages'))
            elif doc_type == 'word':
                return cls._read_word(file_path)
            elif doc_type == 'excel':
                return cls._read_excel(file_path, options.get('excel_sheet'))
            elif doc_type == 'powerpoint':
                return cls._read_powerpoint(file_path, options.get('pptx_slide'))
            else:
                return {
                    'success': False,
                    'error': f'未实现的文档类型: {doc_type}'
                }
        except Exception as e:
            return {
                'success': False,
                'error': f'读取失败: {str(e)}'
            }
    
    @classmethod
    def _read_text(cls, file_path):
        """读取文本文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            'success': True,
            'type': 'text',
            'format': os.path.splitext(file_path)[1].lower(),
            'title': os.path.basename(file_path),
            'content': content,
            'length': len(content)
        }
    
    @classmethod
    def _read_pdf(cls, file_path, pages=None):
        """读取 PDF 文件"""
        try:
            import pdfplumber
        except ImportError:
            return {
                'success': False,
                'error': 'pdfplumber 未安装，请运行: pip install pdfplumber'
            }
        
        text_parts = []
        
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            
            # 解析页面范围
            if pages:
                page_numbers = cls._parse_pages(pages, total_pages)
            else:
                page_numbers = range(1, total_pages + 1)
            
            for page_num in page_numbers:
                if 1 <= page_num <= total_pages:
                    page = pdf.pages[page_num - 1]
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"\n--- 第 {page_num} 页 ---\n")
                        text_parts.append(page_text)
        
        content = '\n'.join(text_parts)
        
        return {
            'success': True,
            'type': 'pdf',
            'format': '.pdf',
            'title': os.path.basename(file_path),
            'content': content,
            'length': len(content),
            'total_pages': total_pages,
            'extracted_pages': len(page_numbers)
        }
    
    @classmethod
    def _read_word(cls, file_path):
        """读取 Word 文档"""
        try:
            from docx import Document
        except ImportError:
            return {
                'success': False,
                'error': 'python-docx 未安装，请运行: pip install python-docx'
            }
        
        doc = Document(file_path)
        
        text_parts = []
        
        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join([cell.text for cell in row.cells])
                text_parts.append(row_text)
        
        content = '\n\n'.join(text_parts)
        
        return {
            'success': True,
            'type': 'word',
            'format': '.docx',
            'title': os.path.basename(file_path),
            'content': content,
            'length': len(content),
            'paragraphs': len(doc.paragraphs),
            'tables': len(doc.tables)
        }
    
    @classmethod
    def _read_excel(cls, file_path, sheet_name=None):
        """读取 Excel 文件"""
        try:
            import openpyxl
        except ImportError:
            return {
                'success': False,
                'error': 'openpyxl 未安装，请运行: pip install openpyxl'
            }
        
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        text_parts = []
        
        # 如果指定了工作表，只读取该表
        if sheet_name and sheet_name in wb.sheetnames:
            sheets = [wb[sheet_name]]
        else:
            sheets = wb.worksheets
        
        for sheet in sheets:
            text_parts.append(f"\n=== 工作表: {sheet.title} ===\n")
            
            # 读取所有单元格
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        content = '\n'.join(text_parts)
        
        return {
            'success': True,
            'type': 'excel',
            'format': '.xlsx',
            'title': os.path.basename(file_path),
            'content': content,
            'length': len(content),
            'sheets': wb.sheetnames,
            'total_sheets': len(wb.sheetnames)
        }
    
    @classmethod
    def _read_powerpoint(cls, file_path, slides=None):
        """读取 PowerPoint 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            return {
                'success': False,
                'error': 'python-pptx 未安装，请运行: pip install python-pptx'
            }
        
        prs = Presentation(file_path)
        
        text_parts = []
        
        # 解析幻灯片范围
        total_slides = len(prs.slides)
        if slides:
            slide_numbers = cls._parse_pages(slides, total_slides)
        else:
            slide_numbers = range(1, total_slides + 1)
        
        for idx, slide in enumerate(prs.slides, 1):
            if idx in slide_numbers:
                text_parts.append(f"\n--- 幻灯片 {idx} ---\n")
                
                # 提取所有形状中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
        
        content = '\n\n'.join(text_parts)
        
        return {
            'success': True,
            'type': 'powerpoint',
            'format': '.pptx',
            'title': os.path.basename(file_path),
            'content': content,
            'length': len(content),
            'total_slides': total_slides,
            'extracted_slides': len(slide_numbers)
        }
    
    @classmethod
    def _parse_pages(cls, pages_str, total):
        """解析页面范围字符串"""
        pages = set()
        
        for part in pages_str.split(','):
            part = part.strip()
            if '-' in part:
                start, end = part.split('-')
                pages.update(range(int(start), int(end) + 1))
            else:
                pages.add(int(part))
        
        # 过滤有效页面
        return [p for p in sorted(pages) if 1 <= p <= total]


# 便捷函数
def read_document(file_path, **options):
    """便捷函数：读取文档"""
    return DocumentReader.read(file_path, **options)


def can_read(file_path):
    """便捷函数：检查是否支持"""
    return DocumentReader.can_read(file_path)


if __name__ == '__main__':
    import json
    
    # 测试
    print("支持的格式:", list(DocumentReader.SUPPORTED_FORMATS.keys()))
    
    # 测试文本文件
    test_file = r'C:\Users\luzhe\.openclaw\workspace-main\skills\summarize\SKILL.md'
    if os.path.exists(test_file):
        result = DocumentReader.read(test_file)
        print(f"\n测试读取文本文件:")
        print(f"  成功: {result['success']}")
        print(f"  类型: {result.get('type')}")
        print(f"  长度: {result.get('length', 0)} 字符")
